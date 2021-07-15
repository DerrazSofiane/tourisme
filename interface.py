"""
Interface graphique pour observer les resultats obtenu du programme main.py
"""

from main import (traitements_informations, generique_variation_valeur_brutes,
                    generique_potentiel, moyenne_donnees_brute_pays, evolutions_sum_annees,
                    tableau_top, tableau_top_pays_hebdo, tableau_top_pays_mensuel,
                    tableau_top_pays_trimestre, comparaison_brute_mois_n)

import seaborn as sns
import matplotlib.pyplot as plt
import streamlit as st
import pandas as pd
import numpy as np
import datetime
from calendar import monthrange
import os
from pptx import Presentation


st.set_option('deprecation.showPyplotGlobalUse', False)

st.title("OBSERVATOIRE DU TOURISME SUR INTERNET")
st.image("https://nicolasbaudy.files.wordpress.com/2020/02/cropped-logo-new-2.png", use_column_width=False)


# Sélection du mode
st.sidebar.write("# Bienvenue :computer:")
mode = st.sidebar.selectbox(
    "1- ANALYSE : choisissez le mode => Générique // Par Pays ",
    ("Générique", "Par pays")
)

# MODE GENERIQUE
st.sidebar.success(f"Vous avez choisi le mode {mode}")
if mode == "Générique":

    try:
        os.mkdir("images_generique")
    except:
        pass
    
    my_expander = st.sidebar.beta_expander("Ouvrir", expanded=True)
    with my_expander:
        st.write("2- DONNÉES : chargez le fichier à analyser :open_file_folder:")
        uploaded_file = st.file_uploader("")

    try:
        fichier = traitements_informations(uploaded_file)
        # Affiche l'intituler de l'objet 
        nom_objet = list(fichier.columns)[0]
        st.sidebar.write(f"Vous analysez: **{nom_objet}** :heavy_check_mark:")
        variation, valeurs_brutes = generique_variation_valeur_brutes(fichier)
        recapitulatif = generique_potentiel(variation, valeurs_brutes)
        top_3 = tableau_top(recapitulatif)


# CALCUL GENERIQUE
        if st.sidebar.checkbox("Calcul Générique") and uploaded_file != "None":
            st.title("Top Potentiel")
            st.write(top_3)
            st.write(top_3.style.set_precision(2))
            st.title("Moyenne des valeurs brutes des 2 dernières semaines")
            st.write(valeurs_brutes.set_index(list(fichier.columns)[0]).style.set_precision(2))
            st.title("Moyenne des variations des 4 dernières semaines")
            st.write(variation.style.set_precision(2))

            
            st.title("Récapitulatif :white_check_mark:")
            st.write(recapitulatif.T.style.set_precision(2))
            
            

# GRAPHIQUE GENERIQUE
        if st.sidebar.checkbox("Graphique Générique") and uploaded_file != "None":
            if st.sidebar.checkbox("Valeurs brutes"):
                tableau = valeurs_brutes
                tableau = tableau.rename({list(fichier.columns)[0]: "semaine"}, axis=1)
                data_melted = pd.melt(tableau, id_vars="semaine", var_name="pays", value_name="valeur")
                st.title("Valeurs brutes de la semaine S et de la semaine (S-1)")
                fig, ax = plt.subplots(figsize=(10,10))
                st.write(sns.barplot(x="pays", y="valeur", hue="semaine", data=data_melted))
                ax.grid(axis="x")
                for p in ax.patches:
                    ax.annotate(format(p.get_height(), '.1f'), 
                        (p.get_x() + p.get_width() / 2., p.get_height()), 
                        ha = 'center', va = 'center', 
                        size=9,
                        xytext = (0, 1), 
                        textcoords = 'offset points')
                st.pyplot()
                

            if st.sidebar.checkbox("Evolution en % de S/S-1"):
                evolution = variation.head(2).reset_index()
                evolution = evolution.rename({"index": "semaine"}, axis=1)
                evolution_melted = pd.melt(evolution.sort_index(ascending=False), id_vars="semaine", var_name="pays", value_name="valeur")
                st.title("Evolution en % de S/S-1")
                fig, ax = plt.subplots(figsize=(10,10))
                st.write(sns.barplot(x="pays", y="valeur", hue="semaine", data=evolution_melted))
                ax.grid(axis="x")
                for p in ax.patches:
                    ax.annotate(format(p.get_height(), '.1f'), 
                        (p.get_x() + p.get_width() / 2., p.get_height()), 
                        ha = 'center', va = 'center', 
                        size=9,
                        xytext = (0, 1), 
                        textcoords = 'offset points')

                st.pyplot()


            if st.sidebar.checkbox("Evolution en % de S-1/S-2"):
                evolution_s1 = variation.tail(2).reset_index()
                evolution_s1 = evolution_s1.rename({"index": "semaine"}, axis=1)
                evolution_s1_melted = pd.melt(evolution_s1.sort_index(ascending=False), id_vars="semaine", var_name="pays", value_name="valeur")
                st.title("Evolution en % de S-1/S-2")
                fig, ax = plt.subplots(figsize=(10,10))
                st.write(sns.barplot(x="pays", y="valeur", hue="semaine", data=evolution_s1_melted))
                ax.grid(axis="x")
                for p in ax.patches:
                    ax.annotate(format(p.get_height(), '.1f'), 
                        (p.get_x() + p.get_width() / 2., p.get_height()), 
                        ha = 'center', va = 'center', 
                        size=9,
                        xytext = (0, 1), 
                        textcoords = 'offset points')

                st.pyplot()

            if st.checkbox("Voulez vous mettre un commentaire ?"):
                commentaire_graph_s2 = st.text_area("Emplacement du commentaire", "")
       
            if st.button("générer un power point Générique"):
                prs = Presentation()
                bullet_slide_layout = prs.slide_layouts[1]
                
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = 'TOURISME TEST'
                
                tf = body_shape.text_frame
                tf.text = '1er slide avec top 3'
                
                p = tf.add_paragraph()
                p.text = 'Use _TextFrame.text for first bullet'
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
                p.level = 2
                
                prs.save('test.pptx')
    except:
        pass


# MODE PAR PAYS
elif mode == "Par pays":
    try:
        os.mkdir("images_par_pays")
    except:
        pass
    my_expander = st.sidebar.beta_expander("Ouvrir", expanded=True)
    with my_expander:
        st.write("Choisir un fichier :open_file_folder:")
        uploaded_file = st.file_uploader("")

    try:
        fichier = traitements_informations(uploaded_file)
        colonnes = list(fichier.columns)
        my_time = datetime.datetime.min.time()
        fichier["Semaine"] = fichier["Semaine"].dt.date
        derniere_date = pd.unique(fichier[colonnes[0]])[-1]
        conv_derniere_date = datetime.datetime.combine(derniere_date, my_time)
        # Ressort une date au format datetime.date
        date_calendar = st.sidebar.date_input("sélectionner la date d'analyse", value=conv_derniere_date)

        if st.sidebar.checkbox("1- Les Top") and uploaded_file != "None":
            recapitualitf_2s, recapitualitf_4s, recapitualitf_12s = moyenne_donnees_brute_pays(fichier, date_calendar)
            st.title("Moyenne des données brutes sur les 2 dernières semaines, des 4 dernières semaines, des 12 dernières semaines")

            cols = st.beta_columns(3)
            cols[0].table(recapitualitf_2s.style.set_precision(2))
            cols[1].table(recapitualitf_4s.style.set_precision(2))
            cols[2].table(recapitualitf_12s.style.set_precision(2))

            st.title("TOP 6")

            cols = st.beta_columns(3)
            cols[0].table(recapitualitf_2s.head(6).style.set_precision(2))
            cols[1].table(recapitualitf_4s.head(6).style.set_precision(2))
            cols[2].table(recapitualitf_12s.head(6).style.set_precision(2))
            if st.checkbox("Voulez vous mettre un commentaire ?"):
                commentaire_recapitualitf_desc = st.text_area("Emplacement du commentaire", "")
                st.write(commentaire_recapitualitf_desc.style.set_precision(2))
            
            recapitualitf_2s, recapitualitf_4s, recapitualitf_12s = moyenne_donnees_brute_pays(fichier, date_calendar)

            if st.sidebar.checkbox("HEBDOMADAIRES") and uploaded_file != "None":

                recapitualitf_2s, recapitualitf_4s, recapitualitf_12s = moyenne_donnees_brute_pays(fichier, date_calendar)
                st.title("Top potentiel hebdo")
                top_pays = tableau_top_pays_hebdo(recapitualitf_2s, fichier)
                st.write(top_pays)
                #if st.checkbox("Voulez vous mettre un commentaire ?"):
                    #commentaire_top_hebdo = st.text_area("Emplacement du commentaire", "")


                if st.sidebar.checkbox("Valeurs brutes des 3 dernères années du top 6 hebdo"):
                    annee = date_calendar.year
                    evolution_annee = evolutions_sum_annees(fichier, annee)
                    top_6 = recapitualitf_2s.head(6)
                    pays = list(top_6.index)
                    annees = pd.unique(evolution_annee["annee"])
                    #f, ax = plt.subplots(2,3,figsize=(10,4))
                    for p in pays:
                        st.write(p)
                        annee1 = evolution_annee[p][(evolution_annee["annee"] == "2019")].reset_index().drop("index", axis=1)
                        annee2 = evolution_annee[p][(evolution_annee["annee"] == "2020")].reset_index().drop("index", axis=1)
                        annee3 = evolution_annee[p][(evolution_annee["annee"] == "2021")].reset_index().drop("index", axis=1)
                        last = pd.concat([annee1, annee2, annee3], axis=1)
                        last.columns = [p+" 2019", p+" 2020", p+" 2021"]
                        last.fillna(0, inplace=True)
                        plt.plot(last)
                        plt.legend(last.columns)
                        st.pyplot()
                
            if st.sidebar.checkbox("MENSUEL") and uploaded_file != "None":
                st.title("Top potentiel mensuel")
                
                top_pays_mensuel = tableau_top_pays_mensuel(recapitualitf_4s, fichier)
                st.write(top_pays_mensuel)
                mois = [i for i in range(1,13)]
                mois_str = ["janvier",	"février",	"mars",	"avril", "mai", "juin",	"juillet", "août", "septembre",
                	"octobre", "novembre", "décembre"]
                mode_mois = st.sidebar.selectbox(
                        "Quel mois?",
                        (mois)
                        )
                derniere_3annee = list(pd.unique(fichier["Semaine"].map(lambda x: x.year)))
                derniere_3annee.sort(reverse=True)
                mode_annee = st.sidebar.selectbox(
                        "Quelle annee?",
                        (derniere_3annee)
                        )
                tableau_brute_n = comparaison_brute_mois_n(fichier, int(mode_mois), int(mode_annee))
                st.title(f"Valeurs brutes du mois {mode_mois} l’année {mode_annee} et du mois {mode_mois} de l’année {mode_annee} -1")
                tableau_brute_n = tableau_brute_n.T
                str_annee = [str(i) for i in derniere_3annee]
                derniere_annee_annee1 = tableau_brute_n[str_annee[0:2]]
                
                derniere_annee_melt = pd.melt(derniere_annee_annee1.reset_index(), id_vars="index", var_name="annee", value_name="valeur")
                st.write(sns.barplot(x="index", y="valeur", hue="annee", data=derniere_annee_melt.sort_values(by=["annee"])))
                
                
                st.pyplot()
                st.title(f"Valeurs brutes du mois {mode_mois} l’année {str(int(mode_annee-1))} et du mois {mode_mois} de l’année {str(int(mode_annee-1))} -1")

                derniere_annee_annee2 = tableau_brute_n[[str_annee[0], str_annee[-1]]]
                
                derniere_annee_melt2 = pd.melt(derniere_annee_annee2.reset_index(), id_vars="index", var_name="annee", value_name="valeur")
                st.write(sns.barplot(x="index", y="valeur", hue="annee", data=derniere_annee_melt2.sort_values(by=["annee"])))

                st.pyplot()
                
            if st.sidebar.checkbox("TRIMESTRE") and uploaded_file != "None":
                st.title("Top potentiel trimestre")
                top_pays = tableau_top_pays_trimestre(recapitualitf_12s, fichier)
                st.write(top_pays)
                
    except:
        pass



